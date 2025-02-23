import sys
import streamlit as st
from langchain.chains import ConversationalRetrievalChain
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.llms import LlamaCpp
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain.memory import ConversationBufferMemory
from langchain.document_loaders import PyPDFLoader
from langchain.docstore.document import Document
import os
import tempfile
from docx import Document as DocxDocument
from pptx import Presentation
import re

sys.path.insert(0, r"C:\Users\Lenovo\AppData\Local\Programs\Python\Python39\python.exe")

def initialize_session_state():
    if 'history' not in st.session_state:
        st.session_state['history'] = []  
    
    if 'generated' not in st.session_state:
        st.session_state['generated'] = ["Hello! Ask me anything about 🤗"]
    
    if 'past' not in st.session_state:
        st.session_state['past'] = ["Hey! 👋"]  

def conversation_chat(query, chain, history):
    
    result = chain({"question": query, "chat_history": history})
    history.append((query, result["answer"]))  
    return result["answer"]

def display_chat_history(chain):
    reply_container = st.container()
    container = st.container()

    with container:
        with st.form(key='my_form', clear_on_submit=True):
            user_input = st.text_input("Question:", placeholder="Ask about your PDF, PPT or DOCX", key='input')
            submit_button = st.form_submit_button(label='Send')

        if submit_button and user_input:
            with st.spinner('Generating response...'):
                user_input = preprocess_input(user_input)
                output = conversation_chat(user_input, chain, st.session_state['history'])
            st.session_state['past'].append(user_input)
            st.session_state['generated'].append(output)

    if st.session_state['generated']:
        with reply_container:
            for i in range(len(st.session_state['generated'])):
                st.chat_message("user").text(st.session_state["past"][i])
                st.chat_message("assistant").text(st.session_state["generated"][i])

        # Allow the user to download the conversation
        download_button = st.download_button(
            label="Download Conversation",
            data="\n".join([f"User: {st.session_state['past'][i]}\nAssistant: {st.session_state['generated'][i]}" for i in range(len(st.session_state['generated']))]),
            file_name="conversation.txt",
            mime="text/plain"
        )

def preprocess_input(input_text):
    cleaned_input = re.sub(r'[^\w\s]', '', input_text.lower())  # Remove special characters
    return cleaned_input

def extract_text_from_word(file_path):
    doc = DocxDocument(file_path)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return text

def extract_text_from_ppt(file_path):
    ppt = Presentation(file_path)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return text

def create_documents_from_text(text_list):
    return [Document(page_content=text) for text in text_list]

def create_conversational_chain(vector_store):
    llm = LlamaCpp(
        streaming=True,
        model_path="mistral-7b-instruct-v0.1.Q4_K_M.gguf",
        temperature=0.75,
        top_p=1,
        verbose=True,
        n_ctx=4096
    )

    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)

    chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        chain_type='stuff',
        retriever=vector_store.as_retriever(search_kwargs={"k": 2}),
        memory=memory
    )
    return chain

def main():
    initialize_session_state()
    st.title("Multi-Document AI Chatbot for PDF, DOCX, and PPT Files :books:")

    st.sidebar.title("Document Processing")
    uploaded_files = st.sidebar.file_uploader("Upload files", accept_multiple_files=True)

    if uploaded_files:
        text = []

        for file in uploaded_files:
            file_extension = os.path.splitext(file.name)[1]

            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                temp_file.write(file.read())
                temp_file_path = temp_file.name

            if file_extension == ".pdf":
                loader = PyPDFLoader(temp_file_path)
                documents = loader.load()
                text.extend([doc.page_content for doc in documents])  # Extract content from each document
            elif file_extension == ".docx":
                text.extend(extract_text_from_word(temp_file_path))
            elif file_extension == ".pptx":
                text.extend(extract_text_from_ppt(temp_file_path))

            os.remove(temp_file_path)

        documents = create_documents_from_text(text)

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=10)
        text_chunks = text_splitter.split_documents(documents)

        embeddings = HuggingFaceEmbeddings(model_name="distilbert-base-nli-mean-tokens", model_kwargs={'device': 'cpu'})

        vector_store = FAISS.from_documents(text_chunks, embedding=embeddings)

        chain = create_conversational_chain(vector_store)

        display_chat_history(chain)

if __name__ == "__main__":
    main()
